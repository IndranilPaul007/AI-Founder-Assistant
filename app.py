import streamlit as st
import pandas as pd
import re
import os
import io
import base64
import json
from dotenv import load_dotenv
from openai import OpenAI
from supabase import create_client, Client
from pptx import Presentation

# Load environment variables
load_dotenv()

# ---------------------------------------------------------
# 1. Page Configuration & Session State Initialization
# ---------------------------------------------------------
st.set_page_config(
    page_title="Founder Assistant | AI Copilot",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Global Session Memory
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
if "user_email" not in st.session_state:
    st.session_state.user_email = ""
if "global_startup_name" not in st.session_state:
    st.session_state.global_startup_name = ""
if "global_idea" not in st.session_state:
    st.session_state.global_idea = ""

if "results" not in st.session_state:
    st.session_state.results = {
        "market": None, "business": None, "competitor": None, 
        "fundraising": None, "execution": None
    }

# View States
if "viewing_history" not in st.session_state:
    st.session_state.viewing_history = None
if "viewing_profile" not in st.session_state:
    st.session_state.viewing_profile = False
if "active_chat_mode" not in st.session_state:
    st.session_state.active_chat_mode = False
if "current_chat_id" not in st.session_state:
    st.session_state.current_chat_id = None
if "messages" not in st.session_state:
    st.session_state.messages = []

# Initialize Supabase Client
@st.cache_resource
def init_supabase() -> Client:
    url = os.getenv("SUPABASE_URL")
    key = os.getenv("SUPABASE_KEY")
    if url and key:
        return create_client(url, key)
    return None

supabase = init_supabase()

# ---------------------------------------------------------
# 2. Cinematic CSS & Deep UI Hiding
# ---------------------------------------------------------
st.markdown("""
   <style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700&display=swap');

:root {
    --bg: #0b1220;
    --surface: #111c2e;
    --surface-hover: #16243a;
    --border: #263854;
    --text: #eef4ff;
    --muted: #9caec7;
    --primary: #3b82f6;
    --primary-hover: #2563eb;
    --accent: #22c55e;
}

/* Base */
html, body, [class*="css"], input, textarea, button {
    font-family: "Inter", sans-serif !important;
}

.stApp {
    background:
        radial-gradient(circle at 75% 0%, rgba(59, 130, 246, 0.14), transparent 32%),
        linear-gradient(135deg, #0b1220 0%, #0e1726 55%, #0a1321 100%);
    color: var(--text);
}

[data-testid="stHeader"] {
    background: rgba(11, 18, 32, 0.78);
    backdrop-filter: blur(12px);
}

footer { visibility: hidden; }

/* Sidebar */
section[data-testid="stSidebar"] {
    background: #0d1626;
    border-right: 1px solid rgba(148, 163, 184, 0.14);
}

section[data-testid="stSidebar"] h1 {
    color: #93c5fd !important;
    font-size: 1.55rem !important;
    font-weight: 700 !important;
    letter-spacing: -0.03em;
}

section[data-testid="stSidebar"] hr {
    border-color: rgba(148, 163, 184, 0.16);
}

/* Main heading classes */
.main-title {
    color: var(--text) !important;
    font-size: clamp(2rem, 4vw, 3rem) !important;
    font-weight: 700 !important;
    letter-spacing: -0.04em;
    margin: 0 0 0.35rem !important;
}

.subtitle {
    color: #93c5fd !important;
    border-left: 3px solid var(--primary);
    padding-left: 0.85rem;
    font-size: 1.05rem !important;
    font-weight: 500;
    margin-bottom: 2rem;
}

/* Typography */
h1, h2, h3, h4, p, label, .stMarkdown {
    color: var(--text) !important;
}

p, .stCaption, [data-testid="stWidgetLabel"] p {
    color: var(--muted) !important;
}

/* Inputs */
.stTextInput input,
.stTextArea textarea,
[data-baseweb="select"] > div {
    background: rgba(15, 27, 45, 0.92) !important;
    color: var(--text) !important;
    border: 1px solid var(--border) !important;
    border-radius: 10px !important;
    box-shadow: none !important;
}

.stTextInput input,
.stTextArea textarea {
    padding: 0.8rem 0.9rem !important;
}

.stTextInput input::placeholder,
.stTextArea textarea::placeholder {
    color: #71839f !important;
    opacity: 1;
}

.stTextInput input:focus,
.stTextArea textarea:focus,
[data-baseweb="select"] > div:focus-within {
    border-color: var(--primary) !important;
    box-shadow: 0 0 0 3px rgba(59, 130, 246, 0.18) !important;
}

/* Buttons */
.stButton > button {
    width: 100%;
    min-height: 2.7rem;
    border: 1px solid transparent !important;
    border-radius: 10px !important;
    background: linear-gradient(135deg, #3b82f6, #2563eb) !important;
    color: #ffffff !important;
    font-weight: 600 !important;
    letter-spacing: 0.01em;
    box-shadow: 0 7px 18px rgba(37, 99, 235, 0.22);
    transition: all 0.2s ease;
}

.stButton > button:hover {
    background: linear-gradient(135deg, #60a5fa, #3b82f6) !important;
    transform: translateY(-1px);
    box-shadow: 0 10px 22px rgba(37, 99, 235, 0.3);
}

.stButton > button:active {
    transform: translateY(0);
}

/* Sidebar radio navigation */
section[data-testid="stSidebar"] div[role="radiogroup"] label {
    border-radius: 9px;
    padding: 0.62rem 0.75rem;
    margin: 0.15rem 0;
    transition: background 0.18s ease;
}

section[data-testid="stSidebar"] div[role="radiogroup"] label:hover {
    background: rgba(59, 130, 246, 0.10);
}

section[data-testid="stSidebar"] div[role="radiogroup"] label[aria-checked="true"] {
    background: rgba(59, 130, 246, 0.17);
    border-left: 3px solid var(--primary);
}

section[data-testid="stSidebar"] div[role="radiogroup"] label[aria-checked="true"] p {
    color: #bfdbfe !important;
    font-weight: 600 !important;
}

/* Tabs */
.stTabs [data-baseweb="tab-list"] {
    gap: 0.4rem;
    border-bottom: 1px solid var(--border);
}

.stTabs [data-baseweb="tab"] {
    color: var(--muted) !important;
    padding: 0.7rem 1rem !important;
}

.stTabs [data-baseweb="tab"][aria-selected="true"] {
    color: #bfdbfe !important;
    border-bottom: 2px solid var(--primary) !important;
    font-weight: 600 !important;
}

/* Optional reusable content card */
.professional-card {
    background: rgba(17, 28, 46, 0.78);
    border: 1px solid var(--border);
    border-radius: 16px;
    padding: 1.4rem;
    box-shadow: 0 12px 30px rgba(0, 0, 0, 0.15);
}

/* Divider */
.cinematic-divider {
    height: 1px;
    margin: 1.5rem 0;
    background: linear-gradient(90deg, transparent, var(--border), transparent);
}
</style>
""", unsafe_allow_html=True)

st.sidebar.markdown('<div class="profile-section">', unsafe_allow_html=True)
st.sidebar.markdown('<div class="profile-name">Indranil Paul</div>', unsafe_allow_html=True)

st.sidebar.markdown('<div class="profile-action">', unsafe_allow_html=True)
if st.sidebar.button("⚙ Profile & Settings", use_container_width=True):
    pass
st.sidebar.markdown('</div>', unsafe_allow_html=True)

st.sidebar.markdown('<div class="logout-action">', unsafe_allow_html=True)
if st.sidebar.button("Disconnect", use_container_width=True):
    pass
st.sidebar.markdown('</div></div>', unsafe_allow_html=True)

# ---------------------------------------------------------
# 3. Helper Functions, Nemotron SDK & PPT Generator
# ---------------------------------------------------------
def is_valid_input(text):
    if not text or len(text.strip()) < 2: return False
    clean = re.sub(r'[^a-z]', '', text.strip().lower())
    if len(clean) < 3: return True
    if len(clean) > 5 and len(re.findall(r'[aeiouy]', clean)) == 0: return False
    if re.search(r'[^aeiouy]{5,}', clean): return False
    if re.search(r'(.)\1{3,}', clean): return False
    for s in ['asdf', 'qwer', 'zxcv', 'hjkl']:
        if s in clean: return False
    return True

def generate_ai_response(api_key, prompt, system_prompt=""):
    if not api_key:
        st.error("⚠️ CRITICAL: API Key Missing. Enter it in the sidebar or .env file.")
        return None
    try:
        client = OpenAI(api_key=api_key, base_url="https://integrate.api.nvidia.com/v1")
        messages = [{"role": "system", "content": system_prompt}] if system_prompt else []
        messages.append({"role": "user", "content": prompt})
        
        response = client.chat.completions.create(
            model="nvidia/nemotron-3-ultra-550b-a55b", 
            messages=messages, temperature=0.7, max_tokens=2048
        )
        return response.choices[0].message.content
    except Exception as e:
        return None

def create_ppt(report_text, report_title):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = report_title
    slide.placeholders[1].text = "Generated by AI Startup Copilot"
    
    sections = report_text.split("## ")
    for section in sections[1:]:
        lines = section.split('\n')
        slide_header = lines[0].strip()
        slide_body = re.sub(r'[*_`#]', '', '\n'.join(lines[1:]).strip())
        if len(slide_body) > 800: slide_body = slide_body[:800] + "...\n\n[Details truncated]"
            
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_header
        slide.placeholders[1].text_frame.text = slide_body
        
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

def display_report_actions(report_text, report_title, state_key):
    st.markdown(report_text)
    st.markdown("<br>", unsafe_allow_html=True)
    col1, col2, col3 = st.columns([1, 1, 2])
    with col1:
        st.download_button("📄 Download (.md)", data=report_text, file_name=f"{report_title.replace(' ','_')}.md", mime="text/markdown", key=f"dl_md_{state_key}")
    with col2:
        try:
            ppt_data = create_ppt(report_text, report_title)
            st.download_button("📊 Download as PPT", data=ppt_data, file_name=f"{report_title.replace(' ','_')}.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", key=f"dl_ppt_{state_key}")
        except: pass

def auto_save_report(report_type, result):
    if supabase:
        supabase.table("saved_reports").insert({
            "user_email": st.session_state.user_email,
            "report_type": report_type,
            "content": result
        }).execute()

# ---------------------------------------------------------
# 4. Authentication Gateway & URL Refresh Fix
# ---------------------------------------------------------
if "session_user" in st.query_params:
    st.session_state.logged_in = True
    st.session_state.user_email = st.query_params["session_user"]

if not st.session_state.logged_in:
    st.markdown("<br><br><br><br>", unsafe_allow_html=True)
    col1, col2, col3 = st.columns([1, 1.5, 1])
    with col2:
        st.markdown("<div class='main-title' style='text-align: center;'>AI Startup Copilot</div>", unsafe_allow_html=True)
        st.markdown("<div class='subtitle' style='text-align: center; border: none;'>Founder Portal</div>", unsafe_allow_html=True)
        
        if supabase is None:
            st.error("⚠️ Supabase Backend Disconnected. Check Secrets.")
            st.stop()
            
        with st.container():
            tab_login, tab_signup = st.tabs(["Login", "Sign Up"])
            with tab_login:
                login_email = st.text_input("Corporate Email", key="l_email")
                login_pwd = st.text_input("Password", type="password", key="l_pwd")
                if st.button("SECURE LOGIN", type="primary", use_container_width=True):
                    if login_email and login_pwd:
                        try:
                            supabase.auth.sign_in_with_password({"email": login_email, "password": login_pwd})
                            st.session_state.logged_in = True
                            st.session_state.user_email = login_email
                            st.query_params["session_user"] = login_email
                            
                            # Check Profile Onboarding
                            res = supabase.table("user_profiles").select("has_onboarded").eq("user_email", login_email).execute()
                            if not res.data:
                                supabase.table("user_profiles").insert({"user_email": login_email}).execute()
                                st.session_state.viewing_profile = True
                            elif not res.data[0].get('has_onboarded'):
                                st.session_state.viewing_profile = True
                                
                            st.rerun()
                        except: st.error("⚠️ Incorrect email or password.")
            
            with tab_signup:
                signup_email = st.text_input("New Email", key="s_email")
                signup_pwd = st.text_input("Create Password (Min 6 chars)", type="password", key="s_pwd")
                if st.button("CREATE ACCOUNT", type="primary", use_container_width=True):
                    if signup_email and len(signup_pwd) >= 6:
                        try:
                            supabase.auth.sign_up({"email": signup_email, "password": signup_pwd})
                            st.success("✅ Account created! Switch to Login.")
                        except Exception as e: st.error(f"⚠️ Error: {e}")
                    else: st.error("⚠️ Password must be 6+ chars.")
    st.stop()

# ---------------------------------------------------------
# Fetch Profile Data & Build Global AI Memory Prompt
# ---------------------------------------------------------
profile_data = supabase.table("user_profiles").select("*").eq("user_email", st.session_state.user_email).execute().data
user_profile = profile_data[0] if profile_data else {}
user_identity = user_profile.get("full_name") or st.session_state.user_email.split('@')[0].capitalize()
ai_memory_context = user_profile.get("ai_memory", "")

global_system_prompt = f"""
You are an elite Silicon Valley venture partner advising {user_identity}.
[CRITICAL USER CONTEXT & MEMORY]:
{ai_memory_context}

Provide exhaustive, factual data-driven startup analysis tailored specifically to the User Context above.
USE PLAIN TEXT ONLY. DO NOT use markdown code blocks or backticks (`). Format cleanly with bullets.
"""

# ---------------------------------------------------------
# 5. Sidebar Configuration & Navigation
# ---------------------------------------------------------
with st.sidebar:
    st.markdown("<h1>AI Copilot</h1>", unsafe_allow_html=True)
    
    # NEW CHAT BUTTON
    if st.button("➕ New Chat Session", type="primary", use_container_width=True):
        st.session_state.messages = [{"role": "assistant", "content": f"Welcome back, {user_identity}! Let's build something."}]
        st.session_state.current_chat_id = None
        st.session_state.active_chat_mode = True
        st.session_state.viewing_profile = False
        st.session_state.viewing_history = None
        st.rerun()
    
    st.markdown("<br>", unsafe_allow_html=True)
    
    # WORKSPACE TOOLS NAVIGATION
    st.markdown("### 🛠️ Workspace Tools")
    
    def on_tool_change():
        st.session_state.viewing_history = None
        st.session_state.viewing_profile = False
        st.session_state.active_chat_mode = False

    selected_tool = st.radio(
        "Tools",
        ["Market Research", "Business Planning", "Competitor Intel", "Fundraising Prep", "Task Execution", "Strategic Advisor"],
        label_visibility="collapsed",
        on_change=on_tool_change
    )
    
    st.markdown("<div class='cinematic-divider'></div>", unsafe_allow_html=True)
    
    # RECENT HISTORY
    st.markdown("### 🗂️ Recent Library")
    if supabase:
        res = supabase.table("saved_reports").select("*").eq("user_email", st.session_state.user_email).order("created_at", desc=True).limit(10).execute()
        for row in res.data:
            date_str = row['created_at'].split("T")[0]
            icon = "💬" if row['report_type'] == "Advisor Chat" else "📄"
            st.markdown('<div class="history-btn">', unsafe_allow_html=True)
            if st.button(f"{icon} {row['report_type']} ({date_str})", key=f"hist_{row['id']}", use_container_width=True):
                if row['report_type'] == "Advisor Chat":
                    st.session_state.messages = json.loads(row['content'])
                    st.session_state.current_chat_id = row['id']
                    st.session_state.active_chat_mode = True
                    st.session_state.viewing_history = None
                    st.session_state.viewing_profile = False
                else:
                    st.session_state.viewing_history = row
                    st.session_state.active_chat_mode = False
                    st.session_state.viewing_profile = False
                st.rerun()
            st.markdown('</div>', unsafe_allow_html=True)

    st.markdown("<div class='cinematic-divider'></div>", unsafe_allow_html=True)
    
    # PROFILE ICON & SETTINGS
    col_p1, col_p2 = st.columns([1, 3])
    with col_p1:
        if user_profile.get("profile_pic"):
            st.markdown(f'<img src="{user_profile["profile_pic"]}" style="border-radius: 50%; width: 40px; height: 40px; object-fit: cover; border: 2px solid #38BDF8;">', unsafe_allow_html=True)
        else:
            st.markdown("👤")
    with col_p2:
        st.caption(f"**{user_identity}**")
        
    if st.button("⚙️ Profile & Settings", use_container_width=True):
        st.session_state.viewing_profile = True
        st.session_state.active_chat_mode = False
        st.session_state.viewing_history = None
        st.rerun()
        
    env_key = os.getenv("NVIDIA_API_KEY", "")
    api_key = env_key if env_key else st.text_input("API Key", type="password")
        
    if st.button("DISCONNECT (Log Out)"):
        supabase.auth.sign_out()
        st.session_state.clear()
        st.query_params.clear()
        st.rerun()

# ---------------------------------------------------------
# 6. View Routing (Overrides main screen based on state)
# ---------------------------------------------------------

# OVERRIDE 1: PROFILE EDITOR
if st.session_state.viewing_profile:
    st.markdown("<div class='main-title'>👤 Founder Profile</div>", unsafe_allow_html=True)
    st.markdown("<div class='subtitle'>Configure your identity and AI memory</div>", unsafe_allow_html=True)
    
    if st.button("⬅️ Return to Workspace", type="primary"):
        st.session_state.viewing_profile = False
        st.rerun()
        
    colA, colB = st.columns([1, 2])
    with colA:
        if user_profile.get("profile_pic"):
            st.markdown(f'<img src="{user_profile["profile_pic"]}" class="profile-pic">', unsafe_allow_html=True)
        img_file = st.file_uploader("Upload Profile Picture", type=["png", "jpg", "jpeg"])
    
    with colB:
        new_name = st.text_input("Full Name", value=user_profile.get("full_name", ""))
        new_phone = st.text_input("Phone Number", value=user_profile.get("phone_number", ""))
        
        st.markdown("#### 🧠 AI Memory Context")
        st.caption("The AI extracts facts from your chats and uses this data across ALL tools to personalize advice.")
        new_memory = st.text_area("What should the AI remember about you?", value=user_profile.get("ai_memory", ""), height=150)
        
        if st.button("💾 Save Profile Data", type="primary"):
            update_data = {
                "full_name": new_name,
                "phone_number": new_phone,
                "ai_memory": new_memory,
                "has_onboarded": True
            }
            if img_file:
                b64 = base64.b64encode(img_file.read()).decode()
                update_data["profile_pic"] = f"data:image/png;base64,{b64}"
                
            supabase.table("user_profiles").update(update_data).eq("user_email", st.session_state.user_email).execute()
            st.success("✅ Profile Updated!")
            st.session_state.viewing_profile = False
            st.rerun()
    st.stop()

# OVERRIDE 2: ARCHIVE VIEWER
if st.session_state.viewing_history:
    rep = st.session_state.viewing_history
    st.markdown(f"<div class='main-title'>🗂️ Archive: {rep['report_type']}</div>", unsafe_allow_html=True)
    
    if st.button("⬅️ Return to Workspace", type="primary"):
        st.session_state.viewing_history = None
        st.rerun()
        
    st.markdown("<div class='cinematic-divider'></div>", unsafe_allow_html=True)
    display_report_actions(rep['content'], rep['report_type'], f"hist_view_{rep['id']}")
    st.stop()

# OVERRIDE 3: FULL SCREEN ACTIVE CHAT (Continued from Sidebar)
if st.session_state.active_chat_mode:
    st.markdown(f"<div class='main-title'>💬 Strategic Advisor</div>", unsafe_allow_html=True)
    if st.button("⬅️ Close Chat & Return to Tools", type="primary"):
        st.session_state.active_chat_mode = False
        st.rerun()
        
    st.markdown("<div class='cinematic-divider'></div>", unsafe_allow_html=True)
    
    chat_container = st.container(height=500)
    with chat_container:
        for msg in st.session_state.messages:
            with st.chat_message(msg["role"]): st.markdown(msg["content"])

    if user_input := st.chat_input("Ask a strategic question..."):
        with chat_container:
            st.chat_message("user").markdown(user_input)
        st.session_state.messages.append({"role": "user", "content": user_input})
        
        with chat_container:
            with st.spinner("Analyzing against verified business facts..."):
                reply = generate_ai_response(api_key, user_input, system_prompt=global_system_prompt)
                
                if reply:
                    st.chat_message("assistant").markdown(reply)
                    st.session_state.messages.append({"role": "assistant", "content": reply})
                    
                    if supabase:
                        content_json = json.dumps(st.session_state.messages)
                        if st.session_state.current_chat_id:
                            supabase.table("saved_reports").update({"content": content_json}).eq("id", st.session_state.current_chat_id).execute()
                        else:
                            res = supabase.table("saved_reports").insert({
                                "user_email": st.session_state.user_email,
                                "report_type": "Advisor Chat",
                                "content": content_json
                            }).execute()
                            st.session_state.current_chat_id = res.data[0]['id']
                            
                        mem_prompt = f"Extract a 1-sentence factual detail about the user's startup or preferences from this message: '{user_input}'. If none, output 'NONE'."
                        new_mem = generate_ai_response(api_key, user_input, system_prompt=mem_prompt)
                        if new_mem and "NONE" not in new_mem and len(new_mem) > 5:
                            updated_mem = ai_memory_context + "\n- " + new_mem.strip()
                            supabase.table("user_profiles").update({"ai_memory": updated_mem}).eq("user_email", st.session_state.user_email).execute()
    st.stop()


# =========================================================
# STANDARD ACTIVE WORKSPACE (Based on Sidebar Selection)
# =========================================================
st.markdown("<div class='main-title'>AI Startup Copilot</div>", unsafe_allow_html=True)
st.markdown(f"<div class='subtitle'>{selected_tool}</div>", unsafe_allow_html=True)

if selected_tool == "Market Research":
    with st.container():
        st.markdown("### Market Intelligence & Deep Validation")
        st.write("Exhaustive analysis: Market Opportunity Score, ICP, TAM/SAM/SOM in ₹ (INR), GTM Channels, Macro Trends, Scenario Matrix, & 48-Hour Experiment.")
        
        col1, col2 = st.columns(2)
        with col1:
            idea = st.text_area("Initialize Core Concept:", value=st.session_state.global_idea, placeholder="Input startup idea or product description...")
            industry = st.text_input("Define Niche / Sector:", placeholder="e.g., Creator Economy, Fintech SaaS, Robotics PaaS")
        with col2:
            geography = st.text_input("Geographic Focus:", placeholder="e.g., India, Global, US/Canada, SE Asia")
            st.markdown("<br>", unsafe_allow_html=True)
            btn_research = st.button("RUN MARKET PROTOCOL", type="primary")

        if btn_research:
            if not is_valid_input(idea):
                st.warning("⚠️ Invalid input detected. Please enter a meaningful startup concept.")
            else:
                st.session_state.global_idea = idea
                with st.spinner("Accessing global market data vectors... Performing deep VC analysis in Rupees (₹)..."):
                    prompt = f"""
                    Perform an exhaustive, highly detailed Venture Capital Market Intelligence Analysis for this startup idea:
                    - Idea: {idea}
                    - Sector: {industry}
                    - Region: {geography}

                    CRITICAL MANDATORY INSTRUCTIONS:
                    - ALL FINANCIAL FIGURES MUST BE PRESENTED STRICTLY IN INDIAN RUPEES (₹ / INR). 
                    - USE PLAIN TEXT ONLY. DO NOT use markdown code blocks, backticks (`), or LaTeX.

                    Structure your response into 8 sections:
                    ## 1. Executive Summary & Market Opportunity Score
                    ## 2. Granular Ideal Customer Profile (ICP) & Buyer Persona
                    ## 3. Comprehensive Market Sizing (TAM / SAM / SOM in ₹)
                    ## 4. Scenario & Sensitivity Analysis Matrix (in ₹)
                    ## 5. Macro Trends & Regulatory Vectors
                    ## 6. Go-To-Market (GTM) Channels & Unit Economics (in ₹)
                    ## 7. Strategic Pros & Hidden Risks
                    ## 8. 48-Hour Tactical Validation Experiment
                    """
                    result = generate_ai_response(api_key, prompt, system_prompt=global_system_prompt)
                    if result:
                        st.session_state.results["market"] = result
                        auto_save_report("Market Research", result)

        if st.session_state.results["market"]:
            st.markdown("#### Deep Market Intelligence Completed:")
            display_report_actions(st.session_state.results["market"], "Market Research", "market")

elif selected_tool == "Business Planning":
    with st.container():
        st.markdown("### Interactive Lean Canvas & Unit Economics (All Financials in ₹)")
        st.write("Generate a structured, operational one-page business model with granular unit economics & a realistic $1B+ unicorn scaling roadmap.")
        
        col1, col2 = st.columns(2)
        with col1:
            startup_name = st.text_input("Startup Identifier (Name):", value=st.session_state.global_startup_name, placeholder="Enter product name...")
            product_desc = st.text_area("Product Description & Core Features:", value=st.session_state.global_idea, placeholder="Describe full product functionality...")
            
            revenue_models_list = [
                "Subscription (SaaS)", "Transactional / Per-Transaction Fee", "Freemium to Enterprise (Upsell)",
                "Usage-Based (Pay-As-You-Go)", "Marketplace Take Rate / Commission", "Direct Sales / E-commerce",
                "Ad-Based / Sponsorship", "Affiliate / Lead Generation", "Open Source / Dual Licensing",
                "Licensing / White-label", "Hardware + Software (Razors & Blades)", "Data Monetization",
                "Franchise Model", "Agency / Retainer / Services", "Other (Specify Below)"
            ]
            selected_model = st.selectbox("Revenue Model:", revenue_models_list)
            
            if selected_model == "Other (Specify Below)":
                custom_model = st.text_input("Specify Custom Revenue Model:", placeholder="e.g., Tokenomic Staking Fee")
                revenue_model = custom_model if custom_model else "Custom Revenue Model"
            else:
                revenue_model = selected_model
            
        with col2:
            target_country = st.text_input("Target Country / Region:", value="India", placeholder="e.g., India, Global")
            value_prop = st.text_area("Core Differentiator (UVP):", placeholder="What makes you uniquely positioned to solve the problem?")
            st.markdown("<br>", unsafe_allow_html=True)
            btn_plan = st.button("BUILD BUSINESS MODEL", type="primary")

    if btn_plan:
        if not is_valid_input(startup_name) or not is_valid_input(product_desc):
            st.warning("⚠️ Invalid input detected. Please provide a valid startup name and product description.")
        else:
            st.session_state.global_startup_name = startup_name
            st.session_state.global_idea = product_desc
            with st.spinner("Synthesizing operational architecture, financial model in ₹, & unicorn scaling roadmap..."):
                prompt = f"""
                Create an exhaustive, professional Lean Business Canvas, Unit Economics blueprint in ₹, and a Realistic Scale Roadmap for:
                - Startup Name: {startup_name}
                - Product Description: {product_desc}
                - Target Region: {target_country}
                - Core Differentiator (UVP): {value_prop}
                - Revenue Model: {revenue_model}

                CRITICAL MANDATORY INSTRUCTIONS:
                - ALL FINANCIAL FIGURES MUST BE PRESENTED STRICTLY IN INDIAN RUPEES (₹ / INR). 
                - USE PLAIN TEXT ONLY. DO NOT use markdown code blocks, backticks (`), or LaTeX.

                Output the analysis in the following sections:
                ## 1. Comprehensive Lean Business Canvas
                ## 2. Granular Unit Economics & Financial Benchmarks (in ₹)
                ## 3. Pragmatic Step-by-Step Scaling Roadmap
                """
                result = generate_ai_response(api_key, prompt, system_prompt=global_system_prompt)
                if result:
                    st.session_state.results["business"] = result
                    auto_save_report("Business Plan", result)
                    
    if st.session_state.results["business"]:
        display_report_actions(st.session_state.results["business"], "Business Plan", "business")

elif selected_tool == "Competitor Intel":
    with st.container():
        st.markdown("### Competitive Intelligence Matrix & Moat Analysis")
        st.write("Exhaustive competitive benchmarking. If target competitors are left blank, AI will auto-discover at least 20 main competitors.")
        
        col1, col2 = st.columns(2)
        with col1:
            my_startup = st.text_input("Your Product Name:", value=st.session_state.global_startup_name, key="c_startup")
            competitors = st.text_input("Target Competitors (Optional):", placeholder="Leave blank for auto-discovery")
        with col2:
            differentiator = st.text_area("Your Defensive Moat / Technology:", placeholder="Why can't they just copy you?")
            st.markdown("<br>", unsafe_allow_html=True)
            btn_comp = st.button("EXECUTE COMPETITIVE BENCHMARK", type="primary")

    if btn_comp:
        if not is_valid_input(my_startup):
            st.warning("⚠️ Please provide a valid product name.")
        else:
            with st.spinner("Analyzing competitive landscape & auto-discovering top 20 competitors..."):
                prompt = f"""
                Perform an exhaustive, institutional-grade competitive intelligence analysis for:
                - Target Startup: {my_startup}
                - Competitors: {competitors if competitors.strip() else 'AUTO-DISCOVER AT LEAST 20 MAIN COMPETITORS GLOBALLY AND IN INDIA.'}
                - Defensive Advantage: {differentiator}

                CRITICAL MANDATORY INSTRUCTIONS:
                - YOU MUST IDENTIFY, NAME, AND ANALYZE AT LEAST 20 MAIN COMPETITORS.
                - ALL FINANCIALS AND PRICING MUST BE IN INDIAN RUPEES (₹ / INR).
                - USE PLAIN TEXT ONLY. DO NOT use markdown code blocks, backticks (`), or LaTeX.

                Output Structure:
                ## 1. Master Competitor Directory (At least 20 Competitors)
                ## 2. Summary Matrix Table
                ## 3. Strategic Blind Spots & Incumbent Vulnerabilities
                ## 4. Defensible Positioning & Moat Lock-in Plan
                """
                result = generate_ai_response(api_key, prompt, system_prompt=global_system_prompt)
                if result:
                    st.session_state.results["competitor"] = result
                    auto_save_report("Competitor Intel", result)
                    
    if st.session_state.results["competitor"]:
        display_report_actions(st.session_state.results["competitor"], "Competitor Intel", "competitor")

elif selected_tool == "Fundraising Prep":
    with st.container():
        st.markdown("### Advanced Investor Protocol & Pitch Prep (Financials in ₹)")
        st.write("Institutional-grade fundraising prep: 30-sec hook, valuation caps in ₹, Cap Table dilution models, term sheet tactics, data room checklist, & VC Q&A.")
        
        col1, col2 = st.columns(2)
        with col1:
            stage = st.selectbox("Current Round Stage:", ["Pre-Seed", "Seed", "Series A", "Series B", "Strategic Partner Round"])
            ask_amount = st.text_input("Target Raise (in ₹):", placeholder="e.g., ₹1.5 Crore or ₹10 Crore")
        with col2:
            traction = st.text_area("Key Milestones / Traction:", placeholder="e.g., 4,500 waitlist, ₹12 Lakhs ARR, 3 enterprise LOIs.")
            st.markdown("<br>", unsafe_allow_html=True)
            btn_fund = st.button("INITIALIZE PITCH ASSETS", type="primary")

    if btn_fund:
        with st.spinner("Generating institutional fundraising suite in ₹..."):
            prompt = f"""
            Generate an advanced, highly detailed investor fundraising suite for:
            - Stage: {stage}
            - Target Raise: {ask_amount}
            - Traction: {traction}

            CRITICAL MANDATORY INSTRUCTIONS: 
            - ALL FINANCIAL FIGURES, VALUATIONS, AND DILUTION MUST BE IN INDIAN RUPEES (₹ / INR). 
            - USE PLAIN TEXT ONLY. DO NOT use markdown code blocks, backticks (`), or LaTeX.

            Provide an advanced 7-part blueprint:
            ## 1. The 30-Second Elevator Pitch & Narrative Hook
            ## 2. Valuation Benchmarks & Cap Table Dilution Scenarios (in ₹)
            ## 3. Investor Persona Targeting Matrix (Angels & VCs in India/Global)
            ## 4. Granular Use of Funds Allocation Breakdown (in ₹)
            ## 5. Term Sheet Negotiation Tactics & Protection Clauses
            ## 6. Investor Due Diligence Data Room Checklist
            ## 7. Top 7 Hardest VC Questions & Battle-Tested Winning Answers
            """
            result = generate_ai_response(api_key, prompt, system_prompt=global_system_prompt)
            if result:
                st.session_state.results["fundraising"] = result
                auto_save_report("Fundraising Prep", result)
                
    if st.session_state.results["fundraising"]:
        display_report_actions(st.session_state.results["fundraising"], "Fundraising Prep", "fundraising")

elif selected_tool == "Task Execution":
    with st.container():
        st.markdown("### Hyper-Realistic 30-Day Execution Roadmap")
        st.write("Convert high-level strategy into metric-driven, day-by-day operational sprints with exact tool stacks and conversion formulas.")
        
        col1, col2 = st.columns(2)
        with col1:
            primary_goal = st.text_input("Critical 30-Day Objective:", placeholder="e.g., Convert 10% of waitlist into paying design partners generating ₹5 Lakhs MRR.")
        with col2:
            weekly_hours = st.slider("Weekly Team Sprints (Hours):", 20, 160, 60)
            st.markdown("<br>", unsafe_allow_html=True)
            btn_task = st.button("GENERATE EXECUTION PLAN", type="primary")

    if btn_task:
        if not is_valid_input(primary_goal):
            st.warning("⚠️ Please provide a meaningful objective.")
        else:
            with st.spinner("Deconstructing strategy into realistic day-by-day operational tasks..."):
                prompt = f"""
                Act as an elite Chief of Staff and Operations Lead. Create a hyper-realistic, granular 30-day execution roadmap to achieve:
                - Goal: {primary_goal}
                - Team Capacity: {weekly_hours} hrs/wk

                CRITICAL MANDATORY INSTRUCTIONS:
                - ALL FINANCIAL METRICS, BUDGETS, AND COSTS MUST BE IN INDIAN RUPEES (₹ / INR).
                - USE PLAIN TEXT ONLY. DO NOT use markdown code blocks, backticks (`), or LaTeX equations. Keep advice highly actionable, practical, and easy to read.

                Structure output into:
                ## 1. Key Success Metrics & Operational Funnel
                ## 2. Recommended Tech Stack & Tooling Suite
                ## 3. Granular Weekly Operational Sprints (Week 1, 2, 3, 4)
                ## 4. Execution Bottlenecks & Operational Contingency Triggers
                """
                result = generate_ai_response(api_key, prompt, system_prompt=global_system_prompt)
                if result:
                    st.session_state.results["execution"] = result
                    auto_save_report("Task Execution", result)
                    
    if st.session_state.results["execution"]:
        display_report_actions(st.session_state.results["execution"], "Task Execution", "execution")

elif selected_tool == "Strategic Advisor":
    with st.container():
        st.markdown("### Strategic Advisory Protocol")
        st.write("Click the button below to initialize a fresh, auto-saving workspace for strategic reasoning.")
        st.markdown("<br>", unsafe_allow_html=True)
        if st.button("🚀 Initialize Live Chat Session", type="primary"):
            st.session_state.messages = [{"role": "assistant", "content": f"Welcome back, {user_identity}! I am your AI Strategic Co-Founder. What critical decision are you evaluating today?"}]
            st.session_state.current_chat_id = None
            st.session_state.active_chat_mode = True
            st.rerun()